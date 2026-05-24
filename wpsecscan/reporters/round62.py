"""Round-62 reporter bundle (#C39-C50).

One module covers many small formats. Larger formats (pptx, docx) get
their own optional-dep wrappers below.

#C39 powerpoint_pptx   — needs python-pptx
#C40 word_docx         — needs python-docx
#C41 jira_bulk_create  — POST every finding to JIRA at once
#C42 confluence_page   — atlassian-python-api (or curl) page upsert
#C43 streamlit_app     — emits a streamlit script the user runs themselves
#C44 grafana_dashboard — JSON dashboard for Grafana datasource
#C45 siem_ndjson       — Splunk HEC / Elastic / Loki / generic SIEM
#C46 datadog_dashboard — JSON for Datadog dashboard API
#C47 csv_pivot         — group findings by check / OWASP / severity
#C48 sbom_diff         — diff two CycloneDX SBOMs
#C49 sbom_vex          — VEX/VDR annotation on a CycloneDX SBOM
#C50 quarterly_trend   — multi-scan trend PDF (uses reportlab if avail)
"""
from __future__ import annotations

import base64
import csv
import io
import json
import os
import time
from collections import Counter
from pathlib import Path
from typing import Any
from urllib.parse import urlparse


# ---- helpers ----

def _summary(report: Any) -> dict:
    if hasattr(report, "to_dict"):
        return report.to_dict()
    return report if isinstance(report, dict) else {}


def _findings_iter(rep: dict):
    for r in rep.get("results", []) or []:
        cid = r.get("check_id", "?")
        for f in r.get("findings", []) or []:
            yield cid, f


# ---- #C39 PowerPoint ----

def powerpoint_pptx(report: Any, out_path: str, *, company_name: str = "") -> str:
    """Generate a 4-slide executive deck. Needs python-pptx (`pip install python-pptx`).

    Returns the path written, or "" on failure.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        return ""
    d = _summary(report)
    s = d.get("summary", {})
    prs = Presentation()

    # Slide 1 — title
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "WPSecScan Executive Report"
    slide.placeholders[1].text = (f"{company_name + ' — ' if company_name else ''}"
                                    f"{d.get('target', '?')}\nGenerated "
                                    f"{time.strftime('%Y-%m-%d %H:%M')}")

    # Slide 2 — risk
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = f"Risk score: {d.get('risk_score', 0)}/100"
    tx = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4)).text_frame
    for sev, label in (("critical", "Critical"), ("high", "High"),
                        ("medium", "Medium"), ("low", "Low"), ("info", "Info")):
        p = tx.add_paragraph()
        p.text = f"{label}: {s.get(sev, 0)}"
        p.font.size = Pt(20)

    # Slide 3 — top critical findings
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Top critical findings"
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5)).text_frame
    count = 0
    for _cid, f in _findings_iter(d):
        if (f.get("severity") or "").lower() != "critical":
            continue
        para = tx.add_paragraph()
        para.text = "• " + (f.get("title") or "")[:120]
        para.font.size = Pt(14)
        count += 1
        if count >= 10:
            break
    if not count:
        tx.add_paragraph().text = "No critical findings."

    # Slide 4 — next steps
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Next steps"
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5)).text_frame
    for step in ("Apply all critical-severity remediations within 24h",
                  "Apply high-severity remediations within 7 days",
                  "Schedule a full re-scan after the fixes",
                  "Brief the team on the new findings"):
        p = tx.add_paragraph()
        p.text = "• " + step
        p.font.size = Pt(16)

    try:
        Path(out_path).parent.mkdir(parents=True, exist_ok=True)
        prs.save(out_path)
        return out_path
    except OSError:
        return ""


# ---- #C40 Word docx ----

def word_docx(report: Any, out_path: str) -> str:
    """Generate a Word .docx report. Needs `python-docx`."""
    try:
        from docx import Document
        from docx.shared import Pt, RGBColor
    except ImportError:
        return ""
    d = _summary(report)
    s = d.get("summary", {})
    doc = Document()
    doc.add_heading("WPSecScan Report", 0)
    doc.add_paragraph(f"Target: {d.get('target', '?')}")
    doc.add_paragraph(f"Generated: {time.strftime('%Y-%m-%d %H:%M')}")
    doc.add_heading(f"Risk score: {d.get('risk_score', 0)}/100", 1)

    tbl = doc.add_table(rows=1, cols=2)
    tbl.style = "Light Grid Accent 1"
    tbl.rows[0].cells[0].text = "Severity"
    tbl.rows[0].cells[1].text = "Count"
    for sev in ("critical", "high", "medium", "low", "info"):
        r = tbl.add_row()
        r.cells[0].text = sev.capitalize()
        r.cells[1].text = str(s.get(sev, 0))

    doc.add_heading("Findings", 1)
    for _cid, f in _findings_iter(d):
        h = doc.add_heading(f"[{(f.get('severity') or 'info').upper()}] {f.get('title', '')}", 3)
        doc.add_paragraph(f"URL: {f.get('url', '')}")
        ev = (f.get("evidence") or "")[:1500]
        if ev:
            p = doc.add_paragraph(ev)
            p.style = doc.styles["Quote"]
        rem = f.get("remediation") or ""
        if rem:
            doc.add_paragraph("Remediation: " + rem[:1500])

    try:
        Path(out_path).parent.mkdir(parents=True, exist_ok=True)
        doc.save(out_path)
        return out_path
    except OSError:
        return ""


# ---- #C41 JIRA bulk create ----

def jira_bulk_create(report: Any, *, project_key: str, severity_min: str = "high",
                       base_url: str | None = None, email: str | None = None,
                       token: str | None = None) -> list[str]:
    """File every finding above `severity_min` as its own JIRA issue.
    Reuses integrations.ticketing.jira_create. Returns list of URLs.
    """
    from ..integrations import ticketing
    base_url = base_url or os.environ.get("JIRA_BASE_URL", "")
    email = email or os.environ.get("JIRA_EMAIL", "")
    token = token or os.environ.get("JIRA_TOKEN", "")
    if not all((base_url, email, token, project_key)):
        return []
    # Temporarily set env for ticketing module
    os.environ["JIRA_BASE_URL"] = base_url
    os.environ["JIRA_EMAIL"] = email
    os.environ["JIRA_TOKEN"] = token

    rank = {"info": 0, "low": 1, "medium": 2, "high": 3, "critical": 4}
    threshold = rank.get(severity_min, 3)
    d = _summary(report)
    target = d.get("target", "")
    urls = []
    for _cid, f in _findings_iter(d):
        if rank.get((f.get("severity") or "").lower(), 0) < threshold:
            continue
        url = ticketing.jira_create(project_key, f, target)
        if url:
            urls.append(url)
    return urls


# ---- #C42 Confluence page export ----

def confluence_page_markdown(report: Any) -> str:
    """Return a Confluence-compatible markdown string. Paste it into
    Confluence's `Insert > Markup > Markdown` panel."""
    d = _summary(report)
    s = d.get("summary", {})
    lines = [
        f"# WPSecScan report — {d.get('target', '?')}",
        f"_Generated {time.strftime('%Y-%m-%d %H:%M')}_",
        "",
        f"**Risk score:** {d.get('risk_score', 0)}/100",
        "",
        "| Severity | Count |",
        "|----------|-------|",
    ]
    for sev in ("critical", "high", "medium", "low", "info"):
        lines.append(f"| {sev.capitalize()} | {s.get(sev, 0)} |")
    lines.append("")
    lines.append("## Findings")
    for cid, f in _findings_iter(d):
        sev = (f.get("severity") or "info").upper()
        lines.append(f"### [{sev}] {f.get('title', '')}")
        lines.append(f"- check: `{cid}`")
        lines.append(f"- url: {f.get('url', '')}")
        ev = (f.get("evidence") or "")[:500]
        if ev:
            lines.append(f"```\n{ev}\n```")
    return "\n".join(lines)


# ---- #C43 Streamlit dashboard ----

def streamlit_script() -> str:
    """Return a self-contained Streamlit script for browsing scan reports.

    Save the output to `dashboard.py` and run:
        pip install streamlit
        streamlit run dashboard.py
    """
    return r"""
# WPSecScan dashboard (generated by reporters/round62.py)
import json
import streamlit as st
from pathlib import Path

st.set_page_config(page_title="WPSecScan dashboard", layout="wide")
st.title("WPSecScan dashboard")

uploaded = st.file_uploader("Upload a WPSecScan JSON report", type=["json"])
if uploaded:
    report = json.loads(uploaded.read())
    st.metric("Risk score", f"{report.get('risk_score', 0)}/100")
    s = report.get("summary", {})
    c1, c2, c3, c4, c5 = st.columns(5)
    c1.metric("Critical", s.get("critical", 0))
    c2.metric("High", s.get("high", 0))
    c3.metric("Medium", s.get("medium", 0))
    c4.metric("Low", s.get("low", 0))
    c5.metric("Info", s.get("info", 0))

    st.divider()
    sev_filter = st.multiselect("Filter by severity",
        ["critical", "high", "medium", "low", "info"], default=["critical", "high"])

    rows = []
    for r in report.get("results", []) or []:
        cid = r.get("check_id", "")
        for f in r.get("findings", []) or []:
            if (f.get("severity") or "") in sev_filter:
                rows.append({
                    "severity": f.get("severity", ""),
                    "title": f.get("title", ""),
                    "check": cid,
                    "url": f.get("url", ""),
                })
    st.dataframe(rows, use_container_width=True, height=600)
else:
    st.info("Drop a wpsecscan-*.json report above to begin.")
"""


# ---- #C44 Grafana dashboard JSON ----

def grafana_dashboard(report: Any) -> dict:
    d = _summary(report)
    s = d.get("summary", {})
    return {
        "title": f"WPSecScan — {d.get('target', '?')}",
        "tags": ["wpsecscan", "security"],
        "panels": [
            {
                "id": 1, "type": "stat",
                "title": f"Risk score / 100",
                "gridPos": {"x": 0, "y": 0, "w": 6, "h": 4},
                "targets": [{"refId": "A", "stat": d.get("risk_score", 0)}],
            },
            {
                "id": 2, "type": "piechart",
                "title": "Severity breakdown",
                "gridPos": {"x": 6, "y": 0, "w": 6, "h": 6},
                "targets": [{"refId": "B", "data": s}],
            },
            {
                "id": 3, "type": "table",
                "title": "Findings",
                "gridPos": {"x": 0, "y": 6, "w": 12, "h": 10},
                "targets": [{
                    "refId": "C",
                    "data": [
                        {"severity": f.get("severity"), "title": f.get("title"),
                          "check": cid, "url": f.get("url")}
                        for cid, f in _findings_iter(d)
                    ],
                }],
            },
        ],
        "schemaVersion": 36,
        "version": 1,
    }


# ---- #C45 SIEM NDJSON (Splunk HEC / Elastic / Loki / generic) ----

def siem_ndjson(report: Any) -> str:
    """Newline-delimited JSON events — one per finding. Ingestable by:
       - Splunk HEC: `curl -k -X POST -d @file.ndjson -H 'Authorization: Splunk <token>' https://splunk:8088/services/collector`
       - Elastic bulk API: prefix each line with `{"index":{"_index":"wpsecscan"}}\n`
       - Loki: wrap with `{"streams":[{"stream":{"app":"wpsecscan"},"values":[...]}]}`
    """
    d = _summary(report)
    target = d.get("target", "")
    lines = []
    for cid, f in _findings_iter(d):
        event = {
            "ts":         int(time.time()),
            "source":     "wpsecscan",
            "target":     target,
            "check_id":   cid,
            "severity":   f.get("severity"),
            "title":      f.get("title"),
            "url":        f.get("url", ""),
            "evidence":   (f.get("evidence") or "")[:2000],
            "remediation": (f.get("remediation") or "")[:2000],
        }
        lines.append(json.dumps(event))
    return "\n".join(lines)


# ---- #C46 Datadog dashboard JSON ----

def datadog_dashboard(report: Any) -> dict:
    d = _summary(report)
    s = d.get("summary", {})
    return {
        "title": f"WPSecScan — {d.get('target', '?')}",
        "description": "Generated by wpsecscan/reporters/round62.py",
        "widgets": [
            {"definition": {"type": "query_value",
                              "title": "Risk score",
                              "requests": [{"q": f"avg:wpsecscan.risk_score{{*}}", "aggregator": "avg"}]}},
            {"definition": {"type": "toplist",
                              "title": "Severity counts",
                              "requests": [{"q": "sum:wpsecscan.findings{*} by {severity}", "limit": 10}]}},
        ],
        "layout_type": "ordered",
        "_summary_seed": s,  # round-trip the data even if it's not native datadog
    }


# ---- #C47 CSV pivot ----

def csv_pivot(report: Any, by: str = "check") -> str:
    """Pivot findings by `check`, `owasp`, or `severity`. Returns CSV string."""
    d = _summary(report)
    counts: Counter = Counter()
    for cid, f in _findings_iter(d):
        key = ""
        if by == "check":
            key = cid
        elif by == "severity":
            key = (f.get("severity") or "info").lower()
        elif by == "owasp":
            # Pull from finding.extra.tags if present
            extra = f.get("extra") or {}
            key = (extra.get("owasp") or "untagged") if isinstance(extra, dict) else "untagged"
        else:
            key = "?"
        counts[key] += 1
    buf = io.StringIO()
    w = csv.writer(buf)
    w.writerow([by, "count"])
    for k, v in sorted(counts.items(), key=lambda kv: -kv[1]):
        w.writerow([k, v])
    return buf.getvalue()


# ---- #C48 SBOM diff ----

def sbom_diff(old_path: str, new_path: str) -> dict:
    """Diff two CycloneDX SBOMs. Returns {added, removed, version_changed}."""
    try:
        old = json.loads(Path(old_path).read_text(encoding="utf-8"))
        new = json.loads(Path(new_path).read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError):
        return {"added": [], "removed": [], "version_changed": []}
    def _components(s: dict) -> dict[str, str]:
        out = {}
        for c in s.get("components", []) or []:
            name = c.get("name") or c.get("bom-ref") or ""
            ver = c.get("version") or ""
            if name:
                out[name] = ver
        return out
    o = _components(old)
    n = _components(new)
    added = sorted(set(n) - set(o))
    removed = sorted(set(o) - set(n))
    changed = sorted(name for name in set(o) & set(n) if o[name] != n[name])
    return {
        "added":            [{"name": x, "version": n[x]} for x in added],
        "removed":          [{"name": x, "version": o[x]} for x in removed],
        "version_changed":  [{"name": x, "old": o[x], "new": n[x]} for x in changed],
    }


# ---- #C49 SBOM VEX/VDR annotation ----

def sbom_vex(report: Any, sbom_path: str, out_path: str) -> str:
    """Walk every CycloneDX component, look up findings that mention it,
    and emit a VEX-annotated SBOM."""
    try:
        sbom = json.loads(Path(sbom_path).read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError):
        return ""
    d = _summary(report)
    # Index findings by mentioned slug
    finding_index: dict[str, list[dict]] = {}
    for cid, f in _findings_iter(d):
        title = (f.get("title") or "").lower()
        for c in sbom.get("components", []) or []:
            name = (c.get("name") or "").lower()
            if name and name in title:
                finding_index.setdefault(name, []).append({"cid": cid, **f})

    # Add vex section to SBOM
    sbom["vex"] = []
    for name, fs in finding_index.items():
        for f in fs:
            sbom["vex"].append({
                "component": name,
                "vulnerability": f.get("title"),
                "severity": f.get("severity"),
                "state": "affected",
                "detail": (f.get("evidence") or "")[:500],
            })
    try:
        Path(out_path).parent.mkdir(parents=True, exist_ok=True)
        Path(out_path).write_text(json.dumps(sbom, indent=2), encoding="utf-8")
        return out_path
    except OSError:
        return ""


# ---- #C50 Quarterly trend PDF ----

def quarterly_trend_pdf(reports: list[Any], out_path: str) -> str:
    """Generate a PDF showing risk-score trend across multiple scans.
    Needs reportlab. Returns path or '' on failure."""
    try:
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.units import cm
    except ImportError:
        return ""
    if not reports:
        return ""
    out = Path(out_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    c = canvas.Canvas(str(out), pagesize=A4)
    w, h = A4
    c.setFont("Helvetica-Bold", 20)
    c.drawString(2 * cm, h - 2 * cm, "WPSecScan quarterly trend")
    c.setFont("Helvetica", 11)
    c.drawString(2 * cm, h - 3 * cm, f"Reports analysed: {len(reports)}")

    # Simple sparkline: dot-and-line of risk_score over time
    pts = []
    for r in reports:
        d = _summary(r)
        pts.append(d.get("risk_score", 0))
    if pts:
        x0, y0 = 2 * cm, h - 7 * cm
        xstep = (w - 4 * cm) / max(1, len(pts) - 1)
        for i in range(len(pts) - 1):
            c.line(x0 + i * xstep, y0 + pts[i],
                   x0 + (i + 1) * xstep, y0 + pts[i + 1])
        # Axis labels
        c.setFont("Helvetica", 9)
        c.drawString(x0, y0 - 0.5 * cm, "earliest")
        c.drawRightString(w - 2 * cm, y0 - 0.5 * cm, "most recent")
        c.drawString(x0 - 1 * cm, y0, "0")
        c.drawString(x0 - 1 * cm, y0 + 100, "100")

    # Summary table
    c.setFont("Helvetica-Bold", 13)
    c.drawString(2 * cm, h - 12 * cm, "Per-scan summary")
    c.setFont("Helvetica", 10)
    y = h - 12.7 * cm
    for i, r in enumerate(reports):
        d = _summary(r)
        s = d.get("summary", {})
        c.drawString(2 * cm, y,
                      f"#{i + 1}  risk={d.get('risk_score', 0)}/100  "
                      f"crit={s.get('critical', 0)}  high={s.get('high', 0)}  "
                      f"med={s.get('medium', 0)}  low={s.get('low', 0)}")
        y -= 0.5 * cm
        if y < 2 * cm:
            c.showPage()
            y = h - 2 * cm
    c.save()
    return str(out)
